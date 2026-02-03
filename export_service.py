"""
Export Service for UDL Differentiation Wizard
Handles exporting materials to DOCX, PDF, PPTX, and XLSX formats
With proper support for checkboxes and bold formatting.
"""

import os
import re
from datetime import datetime
from typing import Optional, List, Tuple

# Document exports
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH

from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_LEFT, TA_CENTER

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN

from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill


VERSION_NAMES = {
    'simplified': 'Simplified (Below Grade Level)',
    'on_level': 'On-Level (Grade Appropriate)',
    'enriched': 'Enriched (Above Grade Level)',
    'visual_heavy': 'Visual-Heavy',
    'scaffolded': 'Step-by-Step Scaffolded'
}

# Unicode checkbox characters
CHECKBOX_UNCHECKED = '\u2610'  # ☐
CHECKBOX_CHECKED = '\u2611'    # ☑


def sanitize_filename(name: str) -> str:
    """Remove invalid characters from filename."""
    return re.sub(r'[<>:"/\\|?*]', '', name)[:50]


def convert_checkboxes(text: str) -> str:
    """Convert markdown-style checkboxes to Unicode checkbox characters."""
    # Convert checked boxes [x] or [X] to ☑
    text = re.sub(r'\[x\]|\[X\]', CHECKBOX_CHECKED, text)
    # Convert unchecked boxes [] or [ ] to ☐
    text = re.sub(r'\[\s?\]', CHECKBOX_UNCHECKED, text)
    return text


def parse_formatted_text(text: str) -> List[Tuple[str, bool, bool]]:
    """
    Parse text with markdown formatting into segments.

    Returns a list of tuples: (text, is_bold, is_italic)
    """
    # First convert checkboxes
    text = convert_checkboxes(text)

    segments = []
    # Pattern to match **bold**, *italic*, or ***bold italic***
    pattern = r'(\*\*\*(.+?)\*\*\*|\*\*(.+?)\*\*|\*(.+?)\*)'

    last_end = 0
    for match in re.finditer(pattern, text):
        # Add any text before this match
        if match.start() > last_end:
            segments.append((text[last_end:match.start()], False, False))

        # Determine formatting
        if match.group(2):  # ***bold italic***
            segments.append((match.group(2), True, True))
        elif match.group(3):  # **bold**
            segments.append((match.group(3), True, False))
        elif match.group(4):  # *italic*
            segments.append((match.group(4), False, True))

        last_end = match.end()

    # Add any remaining text
    if last_end < len(text):
        segments.append((text[last_end:], False, False))

    # If no formatting found, return the whole text
    if not segments:
        segments = [(text, False, False)]

    return segments


def parse_sections(content: str) -> list[tuple[str, str]]:
    """Parse markdown-style sections from content."""
    sections = []
    current_title = "Content"
    current_content = []

    for line in content.split('\n'):
        # Check for markdown headers
        header_match = re.match(r'^#{1,3}\s*\*{0,2}(.+?)\*{0,2}\s*$', line)
        bold_match = re.match(r'^\*\*(.+?)\*\*\s*$', line)

        if header_match or bold_match:
            # Save previous section if it has content
            if current_content:
                sections.append((current_title, '\n'.join(current_content).strip()))
                current_content = []
            current_title = (header_match or bold_match).group(1).strip()
        else:
            current_content.append(line)

    # Add final section
    if current_content:
        sections.append((current_title, '\n'.join(current_content).strip()))

    return sections if sections else [("Content", content)]


def add_formatted_text_to_paragraph(paragraph, text: str):
    """
    Add text with markdown formatting to a DOCX paragraph.
    Handles **bold**, *italic*, and checkbox conversions.
    """
    segments = parse_formatted_text(text)
    for segment_text, is_bold, is_italic in segments:
        run = paragraph.add_run(segment_text)
        if is_bold:
            run.bold = True
        if is_italic:
            run.italic = True


def export_to_docx(materials: dict, form_data: dict, version_key: str,
                   save_path: str) -> str:
    """
    Export a single version to DOCX format.

    Args:
        materials: Dict containing generated materials
        form_data: Form data with learning objective, grade, etc.
        version_key: Which version to export (simplified, on_level, etc.)
        save_path: Directory to save the file

    Returns:
        Full path to saved file
    """
    version_data = materials.get(version_key, {})
    content = version_data.get('content', 'No content generated')
    version_name = VERSION_NAMES.get(version_key, version_key)

    doc = Document()

    # Title
    title = doc.add_heading(f'UDL Learning Materials: {version_name}', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Metadata
    doc.add_paragraph(f"Learning Objective: {form_data.get('learning_objective', 'N/A')}")
    doc.add_paragraph(f"Grade Level: {form_data.get('grade_level', 'N/A')}")
    if form_data.get('subject'):
        doc.add_paragraph(f"Subject: {form_data['subject']}")
    doc.add_paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}")

    doc.add_paragraph()  # Spacer

    # Parse and add content sections
    sections = parse_sections(content)
    for section_title, section_content in sections:
        doc.add_heading(section_title, level=1)

        # Handle bullet points and formatting
        for line in section_content.split('\n'):
            line = line.strip()
            if not line:
                continue

            # Determine the paragraph style and extract content
            if line.startswith('- ') or line.startswith('* '):
                p = doc.add_paragraph(style='List Bullet')
                line_content = line[2:]
            elif re.match(r'^\d+\.\s*', line):
                p = doc.add_paragraph(style='List Number')
                # Remove the number prefix for cleaner formatting
                line_content = re.sub(r'^\d+\.\s*', '', line)
            else:
                p = doc.add_paragraph()
                line_content = line

            # Add the formatted text to the paragraph
            add_formatted_text_to_paragraph(p, line_content)

    # Footer
    doc.add_paragraph()
    footer = doc.add_paragraph()
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_run = footer.add_run("Generated by Assignment Differentiation Wizard")
    footer_run.font.size = Pt(9)
    footer_run.font.italic = True

    # Save
    objective_short = sanitize_filename(form_data.get('learning_objective', 'materials')[:30])
    filename = f"UDL_{version_key}_{objective_short}.docx"
    filepath = os.path.join(save_path, filename)
    doc.save(filepath)

    return filepath


def export_to_pdf(materials: dict, form_data: dict, version_key: str,
                  save_path: str) -> str:
    """Export a single version to PDF format."""
    version_data = materials.get(version_key, {})
    content = version_data.get('content', 'No content generated')
    version_name = VERSION_NAMES.get(version_key, version_key)

    objective_short = sanitize_filename(form_data.get('learning_objective', 'materials')[:30])
    filename = f"UDL_{version_key}_{objective_short}.pdf"
    filepath = os.path.join(save_path, filename)

    doc = SimpleDocTemplate(
        filepath,
        pagesize=letter,
        rightMargin=inch,
        leftMargin=inch,
        topMargin=inch,
        bottomMargin=inch
    )

    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        name='CustomTitle',
        parent=styles['Title'],
        fontSize=18,
        spaceAfter=20,
        alignment=TA_CENTER
    ))
    styles.add(ParagraphStyle(
        name='SectionHeader',
        parent=styles['Heading1'],
        fontSize=14,
        spaceAfter=10,
        spaceBefore=15
    ))
    styles.add(ParagraphStyle(
        name='CustomBody',
        parent=styles['Normal'],
        fontSize=11,
        spaceAfter=6,
        leading=14
    ))

    story = []

    # Title
    story.append(Paragraph(f"UDL Learning Materials: {version_name}", styles['CustomTitle']))
    story.append(Spacer(1, 12))

    # Metadata
    story.append(Paragraph(f"<b>Learning Objective:</b> {form_data.get('learning_objective', 'N/A')}", styles['CustomBody']))
    story.append(Paragraph(f"<b>Grade Level:</b> {form_data.get('grade_level', 'N/A')}", styles['CustomBody']))
    if form_data.get('subject'):
        story.append(Paragraph(f"<b>Subject:</b> {form_data['subject']}", styles['CustomBody']))
    story.append(Paragraph(f"<b>Generated:</b> {datetime.now().strftime('%B %d, %Y')}", styles['CustomBody']))
    story.append(Spacer(1, 20))

    # Content sections
    sections = parse_sections(content)
    for section_title, section_content in sections:
        story.append(Paragraph(section_title, styles['SectionHeader']))

        for line in section_content.split('\n'):
            line = line.strip()
            if line:
                # Convert checkboxes to Unicode characters
                line = convert_checkboxes(line)

                # Convert markdown bold/italic to reportlab HTML tags
                line = re.sub(r'\*\*\*(.+?)\*\*\*', r'<b><i>\1</i></b>', line)
                line = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', line)
                line = re.sub(r'\*(.+?)\*', r'<i>\1</i>', line)

                if line.startswith('- ') or line.startswith('* '):
                    story.append(Paragraph(f"• {line[2:]}", styles['CustomBody']))
                elif re.match(r'^\d+\.\s*', line):
                    story.append(Paragraph(line, styles['CustomBody']))
                else:
                    story.append(Paragraph(line, styles['CustomBody']))

    # Footer
    story.append(Spacer(1, 30))
    story.append(Paragraph(
        "<i>Generated by Assignment Differentiation Wizard</i>",
        ParagraphStyle(name='Footer', parent=styles['Normal'], fontSize=9, alignment=TA_CENTER)
    ))

    doc.build(story)
    return filepath


def export_to_pptx(materials: dict, form_data: dict, version_key: str,
                   save_path: str) -> str:
    """Export a single version to PowerPoint format."""
    version_data = materials.get(version_key, {})
    content = version_data.get('content', 'No content generated')
    version_name = VERSION_NAMES.get(version_key, version_key)

    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = f"UDL Materials: {version_name}"
    subtitle.text = f"{form_data.get('learning_objective', 'Learning Materials')}\n\nGrade: {form_data.get('grade_level', 'N/A')}"

    # Content slides
    sections = parse_sections(content)
    content_layout = prs.slide_layouts[1]  # Title and Content

    for section_title, section_content in sections:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = section_title

        # Add content to body
        body = slide.placeholders[1]
        tf = body.text_frame

        lines = section_content.split('\n')
        first_line = True
        for line in lines:
            line = line.strip()
            if not line:
                continue

            if first_line:
                p = tf.paragraphs[0]
                first_line = False
            else:
                p = tf.add_paragraph()

            # Convert checkboxes first
            line = convert_checkboxes(line)

            # Determine if this is a bullet point
            is_bullet = False
            if line.startswith('- ') or line.startswith('* '):
                line = line[2:]
                is_bullet = True
            elif re.match(r'^\d+\.\s*', line):
                line = re.sub(r'^\d+\.\s*', '', line)
                is_bullet = True

            p.level = 1 if is_bullet else 0

            # Clear default text
            p.text = ""

            # Add formatted runs
            segments = parse_formatted_text(line)
            for segment_text, is_bold, is_italic in segments:
                run = p.add_run()
                run.text = segment_text
                run.font.bold = is_bold
                run.font.italic = is_italic

    # Final slide
    final_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    txBox = final_slide.shapes.add_textbox(PptxInches(1), PptxInches(3), PptxInches(11), PptxInches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Generated by Assignment Differentiation Wizard"
    p.font.size = PptxPt(24)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = datetime.now().strftime('%B %d, %Y')
    p2.font.size = PptxPt(18)
    p2.alignment = PP_ALIGN.CENTER

    # Save
    objective_short = sanitize_filename(form_data.get('learning_objective', 'materials')[:30])
    filename = f"UDL_{version_key}_{objective_short}.pptx"
    filepath = os.path.join(save_path, filename)
    prs.save(filepath)

    return filepath


def export_all_to_xlsx(materials: dict, form_data: dict, save_path: str) -> str:
    """Export all versions to a single Excel workbook."""
    wb = Workbook()

    # Overview sheet
    ws = wb.active
    ws.title = "Overview"

    # Header styling
    header_font = Font(bold=True, size=12)
    header_fill = PatternFill(start_color="6B46C1", end_color="6B46C1", fill_type="solid")
    header_font_white = Font(bold=True, size=12, color="FFFFFF")

    # Overview content
    ws['A1'] = "UDL Differentiation Wizard - Generated Materials"
    ws['A1'].font = Font(bold=True, size=16)
    ws.merge_cells('A1:D1')

    ws['A3'] = "Learning Objective:"
    ws['B3'] = form_data.get('learning_objective', 'N/A')
    ws['A4'] = "Grade Level:"
    ws['B4'] = form_data.get('grade_level', 'N/A')
    ws['A5'] = "Subject:"
    ws['B5'] = form_data.get('subject', 'N/A')
    ws['A6'] = "Generated:"
    ws['B6'] = datetime.now().strftime('%B %d, %Y %H:%M')

    for row in range(3, 7):
        ws[f'A{row}'].font = Font(bold=True)

    ws.column_dimensions['A'].width = 20
    ws.column_dimensions['B'].width = 60

    # Individual version sheets
    for version_key, version_name in VERSION_NAMES.items():
        version_data = materials.get(version_key, {})
        content = version_data.get('content', 'No content generated')

        ws = wb.create_sheet(title=version_key[:31])  # Excel limits sheet names to 31 chars

        ws['A1'] = version_name
        ws['A1'].font = Font(bold=True, size=14)
        ws.merge_cells('A1:C1')

        ws['A3'] = "Section"
        ws['B3'] = "Content"
        ws['A3'].font = header_font_white
        ws['B3'].font = header_font_white
        ws['A3'].fill = header_fill
        ws['B3'].fill = header_fill

        sections = parse_sections(content)
        row = 4
        for section_title, section_content in sections:
            ws[f'A{row}'] = section_title
            ws[f'A{row}'].font = Font(bold=True)
            ws[f'A{row}'].alignment = Alignment(vertical='top')

            # Convert checkboxes and clean up markdown for Excel
            cleaned_content = convert_checkboxes(section_content)
            # Remove markdown bold/italic markers for plain text display
            cleaned_content = re.sub(r'\*\*\*(.+?)\*\*\*', r'\1', cleaned_content)
            cleaned_content = re.sub(r'\*\*(.+?)\*\*', r'\1', cleaned_content)
            cleaned_content = re.sub(r'\*(.+?)\*', r'\1', cleaned_content)

            ws[f'B{row}'] = cleaned_content
            ws[f'B{row}'].alignment = Alignment(wrap_text=True, vertical='top')
            row += 1

        ws.column_dimensions['A'].width = 25
        ws.column_dimensions['B'].width = 80

    # Save
    objective_short = sanitize_filename(form_data.get('learning_objective', 'materials')[:30])
    filename = f"UDL_AllVersions_{objective_short}.xlsx"
    filepath = os.path.join(save_path, filename)
    wb.save(filepath)

    return filepath
